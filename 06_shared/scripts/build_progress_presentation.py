#!/usr/bin/env python3
"""Generate LiSPER progress-report PowerPoint for advisor meeting."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx.chart.data import CategoryChartData
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "05_manuscript" / "LiSPER_Progress_Report.pptx"
BANNER = ROOT / "assets" / "branding" / "banners" / "Dark_Banner.png"

# LiSPER palette — deep navy + electric cyan (matches README branding)
NAVY = RGBColor(0x0B, 0x1A, 0x33)
NAVY_MID = RGBColor(0x12, 0x2B, 0x4F)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
CYAN_DARK = RGBColor(0x06, 0xB6, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def add_slide_transition(slide, effect: str = "fade") -> None:
    """Add a medium-speed slide transition."""
    sld = slide._element
    for old in sld.findall(f"{{{NS_P}}}transition"):
        sld.remove(old)
    transition = etree.Element(f"{{{NS_P}}}transition")
    transition.set("spd", "med")
    child = etree.SubElement(transition, f"{{{NS_P}}}{effect}")
    sld.append(transition)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor, alpha: float | None = None) -> None:
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if alpha is not None:
        shape.fill.transparency = alpha
    shape.line.fill.background()


def add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(1.05), NAVY)
    add_rect(slide, Inches(0), Inches(1.05), Inches(10), Inches(0.06), CYAN)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(8.8), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.62), Inches(8.8), Inches(0.35))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = CYAN
        sp.font.name = "Calibri"


def add_bullets(slide, items: list[str], x, y, w, h, size: int = 16, color: RGBColor = TEXT) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)


def add_status_cards(slide, cards: list[tuple[str, str, str, RGBColor]], x, y, w, h) -> None:
    """cards: (label, status, detail, accent_color)"""
    n = len(cards)
    gap = Inches(0.12)
    card_w = (w - gap * (n - 1)) / n
    for i, (label, status, detail, accent) in enumerate(cards):
        cx = x + i * (card_w + gap)
        add_rect(slide, cx, y, card_w, h, WHITE)
        add_rect(slide, cx, y, card_w, Inches(0.08), accent)
        t1 = slide.shapes.add_textbox(cx + Inches(0.12), y + Inches(0.18), card_w - Inches(0.24), Inches(0.35))
        t1.text_frame.paragraphs[0].text = label
        t1.text_frame.paragraphs[0].font.size = Pt(11)
        t1.text_frame.paragraphs[0].font.color.rgb = SLATE
        t1.text_frame.paragraphs[0].font.bold = True
        t2 = slide.shapes.add_textbox(cx + Inches(0.12), y + Inches(0.5), card_w - Inches(0.24), Inches(0.45))
        t2.text_frame.paragraphs[0].text = status
        t2.text_frame.paragraphs[0].font.size = Pt(20)
        t2.text_frame.paragraphs[0].font.bold = True
        t2.text_frame.paragraphs[0].font.color.rgb = accent
        t3 = slide.shapes.add_textbox(cx + Inches(0.12), y + Inches(1.0), card_w - Inches(0.24), h - Inches(1.15))
        tf = t3.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = detail
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = TEXT


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, NAVY)
    if BANNER.exists():
        s1.shapes.add_picture(str(BANNER), Inches(0), Inches(0), width=Inches(10), height=Inches(2.35))
    add_rect(s1, Inches(0), Inches(2.35), Inches(10), Inches(3.275), NAVY_MID, alpha=0.15)
    t = s1.shapes.add_textbox(Inches(0.6), Inches(2.55), Inches(8.8), Inches(0.8))
    t.text_frame.paragraphs[0].text = "LiSPER Progress Report"
    t.text_frame.paragraphs[0].font.size = Pt(36)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub = s1.shapes.add_textbox(Inches(0.6), Inches(3.25), Inches(8.8), Inches(0.5))
    sub.text_frame.paragraphs[0].text = "Lithium-Selective Peptide Engineering and Recovery"
    sub.text_frame.paragraphs[0].font.size = Pt(16)
    sub.text_frame.paragraphs[0].font.color.rgb = CYAN
    meta = s1.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(8.8), Inches(1.2))
    mtf = meta.text_frame
    for i, line in enumerate(
        [
            "Advisor progress update — computational validation & experimental preparation",
            "Jacky Lin  |  June 2026",
            "github.com/jackyissocute/LiSPER",
        ]
    ):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE
    add_slide_transition(s1, "fade")

    # --- Slide 2: Scientific question ---
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, LIGHT_BG)
    add_title_bar(s2, "Scientific Question", "Can engineered peptides prefer Li+ over Na+?")
    add_bullets(
        s2,
        [
            "LiSPER designs short, IDP-like peptides with oxygen-donor residues (Asp, Ser) and literature GPGDP/GPGNP motifs.",
            "Goal: measurable Li+/Na+ selectivity in aqueous solution — not just Li+ binding.",
            "Long-term vision: peptide-enabled Bio-DLE for battery recycling and lithium recovery streams.",
            "Selectivity metric:  ΔΔG = ΔG(Li+) − ΔG(Na+)  (more negative = stronger Li+ preference).",
        ],
        Inches(0.55), Inches(1.35), Inches(5.6), Inches(3.8),
    )
    add_rect(s2, Inches(6.35), Inches(1.45), Inches(3.1), Inches(3.55), NAVY)
    flow = s2.shapes.add_textbox(Inches(6.55), Inches(1.65), Inches(2.7), Inches(3.2))
    ftf = flow.text_frame
    ftf.word_wrap = True
    steps = [
        ("Motif inspiration", "GPGDP / GPGNP"),
        ("IDP flexibility", "Gly / Ser / Pro rich"),
        ("Paired simulation", "LiCl vs NaCl MD"),
        ("Free-energy ranking", "Umbrella → PMF"),
        ("Wet-lab validation", "His6-SUMO expression"),
    ]
    for i, (head, body) in enumerate(steps):
        p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
        p.text = head
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.space_after = Pt(2)
        p2 = ftf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(11)
        p2.font.color.rgb = WHITE
        p2.space_after = Pt(10)
    add_slide_transition(s2, "push")

    # --- Slide 3: Three-phase program ---
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, LIGHT_BG)
    add_title_bar(s3, "Three-Phase Program", "Discovery → Validation → Translation")
    add_status_cards(
        s3,
        [
            ("Phase I", "ACTIVE", "Computational discovery\n• 10 candidates\n• MD + clustering\n• PMF ranking planned", CYAN_DARK),
            ("Phase II", "PREPARING", "Experimental validation\n• His6-SUMO plasmids ready\n• Expression & assays next", GREEN),
            ("Phase III", "CONCEPT", "Industrial translation\n• Host selection report\n• Packed-bed Bio-DLE architecture", AMBER),
        ],
        Inches(0.55), Inches(1.45), Inches(8.9), Inches(2.0),
    )
    add_bullets(
        s3,
        [
            "Repository reorganized (June 2026) into stage-based folders for clearer tracking.",
            "Professor-relevant focus today: computational pipeline status + vendor-ready plasmid package.",
        ],
        Inches(0.55), Inches(3.75), Inches(8.9), Inches(1.4),
        size=14,
    )
    add_slide_transition(s3, "fade")

    # --- Slide 4: Candidate library ---
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, LIGHT_BG)
    add_title_bar(s4, "First-Round Candidate Library", "10 peptides — compact enough for simulation and wet-lab follow-through")
    rows = [
        ["Rank", "Candidate", "Sequence (abbrev.)", "First subset?"],
        ["1", "LiD3-1", "GPGDP ×3 repeat", "Yes"],
        ["2", "LiND-1", "GPGNP + GPGDP hybrid", "Yes"],
        ["3", "IDP-Li-1", "Flexible acidic shell", "Yes"],
        ["5", "LowCharge-Li", "Lower charge control", "Yes"],
        ["10", "Control-Negative", "Weak/neutral control", "Yes"],
        ["4–9", "4 additional designs", "Diversity / affinity variants", "Screen later"],
    ]
    table = s4.shapes.add_table(len(rows), 4, Inches(0.55), Inches(1.35), Inches(8.9), Inches(3.5)).table
    col_w = [Inches(0.7), Inches(1.6), Inches(3.8), Inches(1.2)]
    for i, w in enumerate(col_w):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if r else 12)
                p.font.name = "Calibri"
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    note = s4.shapes.add_textbox(Inches(0.55), Inches(4.95), Inches(8.9), Inches(0.45))
    note.text_frame.paragraphs[0].text = "Wet-lab first subset: LiD3-1, LiND-1, IDP-Li-1, LowCharge-Li, Control-Negative"
    note.text_frame.paragraphs[0].font.size = Pt(12)
    note.text_frame.paragraphs[0].font.color.rgb = CYAN_DARK
    note.text_frame.paragraphs[0].font.bold = True
    add_slide_transition(s4, "wipe")

    # --- Slide 5: Computational pipeline complete stages ---
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, LIGHT_BG)
    add_title_bar(s5, "Computational Pipeline — Completed Stages", "ESMFold → CHARMM-GUI → Minimization → Equilibration")
    stages = [
        ("Candidate design", "Complete", "10 sequences in candidates.tsv", GREEN),
        ("ESMFold structures", "Complete", "PDB + PAE for all 10; starting models only", GREEN),
        ("CHARMM-GUI systems", "Complete", "10 LiCl + 10 NaCl (CHARMM36m, TIP3P)", GREEN),
        ("GROMACS minimize", "Complete", "10/10 LiCl and NaCl", GREEN),
        ("GROMACS equilibrate", "Complete", "10/10 LiCl and NaCl (step4.1)", GREEN),
    ]
    y = Inches(1.35)
    for label, status, detail, color in stages:
        add_rect(s5, Inches(0.55), y, Inches(0.12), Inches(0.62), color)
        add_rect(s5, Inches(0.75), y, Inches(8.7), Inches(0.62), WHITE)
        t1 = s5.shapes.add_textbox(Inches(0.9), y + Inches(0.08), Inches(2.2), Inches(0.3))
        t1.text_frame.paragraphs[0].text = label
        t1.text_frame.paragraphs[0].font.size = Pt(14)
        t1.text_frame.paragraphs[0].font.bold = True
        t2 = s5.shapes.add_textbox(Inches(3.1), y + Inches(0.08), Inches(1.3), Inches(0.3))
        t2.text_frame.paragraphs[0].text = status
        t2.text_frame.paragraphs[0].font.size = Pt(13)
        t2.text_frame.paragraphs[0].font.color.rgb = color
        t2.text_frame.paragraphs[0].font.bold = True
        t3 = s5.shapes.add_textbox(Inches(4.5), y + Inches(0.08), Inches(4.8), Inches(0.45))
        t3.text_frame.paragraphs[0].text = detail
        t3.text_frame.paragraphs[0].font.size = Pt(12)
        t3.text_frame.paragraphs[0].font.color.rgb = SLATE
        y += Inches(0.72)
    add_bullets(
        s5,
        ["Ensemble-aware workflow: ESMFold is a starting point; selectivity comes from MD → clustering → PMF."],
        Inches(0.55), Inches(5.05), Inches(8.9), Inches(0.45),
        size=12,
        color=SLATE,
    )
    add_slide_transition(s5, "fade")

    # --- Slide 6: Progress dashboard chart ---
    s6 = prs.slides.add_slide(blank)
    set_bg(s6, LIGHT_BG)
    add_title_bar(s6, "Project Dashboard", "Workstream progress from repository status files")
    chart_data = CategoryChartData()
    chart_data.categories = [
        "Candidate design",
        "ESMFold",
        "CHARMM-GUI",
        "Equilibration",
        "Production MD",
        "Umbrella PMF",
        "Plasmids",
        "Wet-lab",
    ]
    chart_data.add_series("Progress (%)", (100, 100, 100, 100, 25, 5, 95, 15))
    chart_frame = s6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.55), Inches(1.35), Inches(5.5), Inches(3.8),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN_DARK
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    val_axis = chart.value_axis
    val_axis.maximum_scale = 100
    val_axis.has_major_gridlines = True
    legend_box = s6.shapes.add_textbox(Inches(6.3), Inches(1.45), Inches(3.2), Inches(3.6))
    ltf = legend_box.text_frame
    ltf.word_wrap = True
    legend_items = [
        ("Complete", "Candidate design, ESMFold, CHARMM-GUI, equilibration (20/20 systems)"),
        ("Active (~25%)", "20 ns production MD + clustering — LiD3-1 done; IDP-Li-1 running"),
        ("Planned (~5%)", "Umbrella sampling and ΔΔG PMF ranking"),
        ("Preparing (~95%)", "His6-SUMO plasmid package vendor-ready; wet-lab not started"),
    ]
    for i, (head, body) in enumerate(legend_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = head
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_after = Pt(2)
        p2 = ltf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT
        p2.space_after = Pt(12)
    add_slide_transition(s6, "push")

    # --- Slide 7: MD production status ---
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, LIGHT_BG)
    add_title_bar(s7, "MD Production Status (LiCl)", "Remote GROMACS on SeeTACloud — synced 2026-06-15")
    md_rows = [
        ["Candidate", "Production", "Clustering", "Notes"],
        ["LiD3-1", "✓ 20 ns complete", "Blocked (trjconv)", "Trajectory valid; index/atom mismatch at centering"],
        ["LiND-1", "Blocked (grompp)", "—", "Missing toppar/forcefield.itp in prod context"],
        ["IDP-Li-1", "~43% (8.65/20 ns)", "Pending", "Active run; T ≈ 295 K, stable constraints"],
        ["Remaining 7", "Queued", "Pending", "Sequential CPU-only queue after IDP-Li-1"],
        ["NaCl (all 10)", "Waiting", "Pending", "Queued behind LiCl production/clustering"],
    ]
    t7 = s7.shapes.add_table(len(md_rows), 4, Inches(0.55), Inches(1.35), Inches(8.9), Inches(2.8)).table
    widths = [Inches(1.4), Inches(1.8), Inches(1.5), Inches(4.2)]
    for i, w in enumerate(widths):
        t7.columns[i].width = w
    for r, row in enumerate(md_rows):
        for c, val in enumerate(row):
            cell = t7.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10 if r else 11)
                p.font.name = "Calibri"
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
    add_bullets(
        s7,
        [
            "Scientific gate: obtain clustered representative structures before umbrella sampling.",
            "LiD3-1 production succeeded — clustering repair is a post-processing fix, not failed MD.",
            "Rough estimate: ~8–10 days remaining LiCl queue + ~10 days NaCl (CPU-only, sequential).",
        ],
        Inches(0.55), Inches(4.35), Inches(8.9), Inches(1.1),
        size=12,
    )
    add_slide_transition(s7, "fade")

    # --- Slide 8: MD workflow diagram ---
    s8 = prs.slides.add_slide(blank)
    set_bg(s8, LIGHT_BG)
    add_title_bar(s8, "Computational Validation Workflow", "Why clustering matters for IDP-like peptides")
    boxes = [
        ("ESMFold", "Starting structure"),
        ("CHARMM-GUI", "Solvated LiCl / NaCl"),
        ("20 ns MD", "Conformational ensemble"),
        ("Cluster", "Top-populated state"),
        ("Umbrella", "PMF windows"),
        ("ΔΔG rank", "Li+ vs Na+ selectivity"),
    ]
    bx = Inches(0.4)
    bw = Inches(1.45)
    bh = Inches(0.95)
    by = Inches(2.0)
    for i, (title, sub) in enumerate(boxes):
        cx = bx + i * (bw + Inches(0.08))
        color = GREEN if i < 3 else (CYAN_DARK if i == 3 else (AMBER if i < 5 else PURPLE))
        add_rect(s8, cx, by, bw, bh, color if i >= 3 else NAVY)
        tb = s8.shapes.add_textbox(cx + Inches(0.08), by + Inches(0.12), bw - Inches(0.16), bh - Inches(0.2))
        tf = tb.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(9)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        if i < len(boxes) - 1:
            arr = s8.shapes.add_textbox(cx + bw, by + Inches(0.35), Inches(0.12), Inches(0.3))
            arr.text_frame.paragraphs[0].text = "→"
            arr.text_frame.paragraphs[0].font.size = Pt(18)
            arr.text_frame.paragraphs[0].font.color.rgb = SLATE
    add_bullets(
        s8,
        [
            "Force field: CHARMM36m + TIP3P  |  GROMACS 2026.0  |  Remote CPU-only runs",
            "We are currently between steps 3 and 4 — first complete trajectory exists (LiD3-1 LiCl).",
        ],
        Inches(0.55), Inches(3.35), Inches(8.9), Inches(1.5),
        size=13,
    )
    add_slide_transition(s8, "wipe")

    # --- Slide 9: Plasmid rationale ---
    s9 = prs.slides.add_slide(blank)
    set_bg(s9, LIGHT_BG)
    add_title_bar(s9, "Experimental Validation — Plasmid Design", "His6-SUMO fusion replaces deprecated direct His/T7 constructs")
    add_bullets(
        s9,
        [
            "Problem with v1 designs: peptides too small (~1.1–1.6 kDa) for robust expression; vector-derived residues on assay molecule.",
            "Solution: express as His6-Smt3 SUMO fusion (~13.3–13.8 kDa) in pET-28a(+), E. coli BL21(DE3).",
            "SUMO protease cleavage after C-terminal GG releases exact native LiSPER peptide — no His-tag remnant.",
            "All 10 candidates have vendor-ready GenBank maps, FASTA inserts, and order tables.",
        ],
        Inches(0.55), Inches(1.35), Inches(5.5), Inches(3.5),
    )
    arch = s9.shapes.add_textbox(Inches(6.2), Inches(1.55), Inches(3.25), Inches(1.0))
    arch.text_frame.paragraphs[0].text = "Expression architecture"
    arch.text_frame.paragraphs[0].font.size = Pt(12)
    arch.text_frame.paragraphs[0].font.bold = True
    arch.text_frame.paragraphs[0].font.color.rgb = NAVY
    add_rect(s9, Inches(6.2), Inches(2.0), Inches(3.25), Inches(0.55), NAVY)
    cas = s9.shapes.add_textbox(Inches(6.3), Inches(2.08), Inches(3.05), Inches(0.4))
    cas.text_frame.paragraphs[0].text = "T7 → RBS → His6 → SUMO → LiSPER → STOP"
    cas.text_frame.paragraphs[0].font.size = Pt(11)
    cas.text_frame.paragraphs[0].font.color.rgb = CYAN
    cas.text_frame.paragraphs[0].font.bold = True
    add_rect(s9, Inches(6.2), Inches(2.75), Inches(3.25), Inches(1.85), WHITE)
    pkg = s9.shapes.add_textbox(Inches(6.35), Inches(2.9), Inches(2.95), Inches(1.6))
    ptf = pkg.text_frame
    ptf.word_wrap = True
    for i, line in enumerate(
        [
            "vendor_ready_SUMO/",
            "• 10 GenBank plasmid maps",
            "• Synthetic ORF FASTA files",
            "• Construct summary CSV",
            "• Vendor submission report",
            "• Reproducible Python generator",
        ]
    ):
        p = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT if i else CYAN_DARK
        p.font.bold = i == 0
    add_slide_transition(s9, "fade")

    # --- Slide 10: Wet-lab plan ---
    s10 = prs.slides.add_slide(blank)
    set_bg(s10, LIGHT_BG)
    add_title_bar(s10, "Wet-Lab Plan (Phase II)", "Ready to order plasmids — assays designed, not yet executed")
    add_status_cards(
        s10,
        [
            ("Expression", "Planned", "BL21(DE3), T7 induction, Kan selection\nTris-Tricine SDS-PAGE (~14 kDa fusion)", AMBER),
            ("Purification", "Planned", "Ni-NTA native conditions\nSUMO protease cleavage", AMBER),
            ("Peptide recovery", "Planned", "Second Ni-NTA → untagged peptide in flow-through\nLC-MS / MALDI identity check", AMBER),
            ("Binding assays", "Planned", "Li+ binding + Na+ competition\nFirst subset: 5 candidates", GREEN),
        ],
        Inches(0.55), Inches(1.35), Inches(8.9), Inches(2.15),
    )
    add_bullets(
        s10,
        [
            "Peptides lack aromatic residues — A280 is not suitable; use mass spec for identity.",
            "Low-MW peptides (~1 kDa) will not be retained on standard 3–30 kDa MWCO devices.",
            "Next experimental gate: vendor synthesis → sequence verification → pilot expression of first subset.",
        ],
        Inches(0.55), Inches(3.75), Inches(8.9), Inches(1.5),
        size=13,
    )
    add_slide_transition(s10, "push")

    # --- Slide 11: Achievements summary ---
    s11 = prs.slides.add_slide(blank)
    set_bg(s11, NAVY)
    add_rect(s11, Inches(0), Inches(0), Inches(10), Inches(0.08), CYAN)
    h = s11.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.6))
    h.text_frame.paragraphs[0].text = "Key Achievements to Date"
    h.text_frame.paragraphs[0].font.size = Pt(30)
    h.text_frame.paragraphs[0].font.bold = True
    h.text_frame.paragraphs[0].font.color.rgb = WHITE
    achievements = [
        ("10", "LiSPER candidates designed with literature + IDP rationale"),
        ("20", "Solvated MD systems built (10 LiCl + 10 NaCl)"),
        ("20/20", "Systems minimized & equilibrated"),
        ("1", "Complete 20 ns LiCl production trajectory (LiD3-1)"),
        ("10", "Vendor-ready His6-SUMO plasmid constructs"),
    ]
    for i, (num, label) in enumerate(achievements):
        col = i % 3
        row = i // 3
        cx = Inches(0.55) + col * Inches(3.05)
        cy = Inches(1.2) + row * Inches(1.85)
        add_rect(s11, cx, cy, Inches(2.85), Inches(1.55), NAVY_MID)
        nbox = s11.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.15), Inches(2.55), Inches(0.7))
        nbox.text_frame.paragraphs[0].text = num
        nbox.text_frame.paragraphs[0].font.size = Pt(36)
        nbox.text_frame.paragraphs[0].font.bold = True
        nbox.text_frame.paragraphs[0].font.color.rgb = CYAN
        lbox = s11.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.85), Inches(2.55), Inches(0.6))
        lbox.text_frame.word_wrap = True
        lbox.text_frame.paragraphs[0].text = label
        lbox.text_frame.paragraphs[0].font.size = Pt(11)
        lbox.text_frame.paragraphs[0].font.color.rgb = WHITE
    add_slide_transition(s11, "fade")

    # --- Slide 12: Next steps & timeline ---
    s12 = prs.slides.add_slide(blank)
    set_bg(s12, LIGHT_BG)
    add_title_bar(s12, "Next Steps & Timeline", "Parallel computational queue + experimental kickoff")
    timeline = [
        ("Now – 2 weeks", "Finish LiCl/NaCl production MD; repair LiD3-1 clustering & LiND-1 grompp", CYAN_DARK),
        ("Weeks 2–4", "Structural clustering → representative PDBs for all candidates", PURPLE),
        ("Weeks 4–8", "Umbrella sampling + PMF → ΔΔG Li+/Na+ ranking", PURPLE),
        ("Immediate", "Submit His6-SUMO plasmids to vendor; begin with 5-candidate subset", GREEN),
        ("After plasmids arrive", "Expression → Ni-NTA → SUMO cleavage → Li+/Na+ binding assays", GREEN),
    ]
    y = Inches(1.35)
    for period, task, color in timeline:
        add_rect(s12, Inches(0.55), y, Inches(1.65), Inches(0.58), color)
        p1 = s12.shapes.add_textbox(Inches(0.62), y + Inches(0.1), Inches(1.5), Inches(0.4))
        p1.text_frame.paragraphs[0].text = period
        p1.text_frame.paragraphs[0].font.size = Pt(10)
        p1.text_frame.paragraphs[0].font.bold = True
        p1.text_frame.paragraphs[0].font.color.rgb = WHITE
        add_rect(s12, Inches(2.35), y, Inches(7.1), Inches(0.58), WHITE)
        p2 = s12.shapes.add_textbox(Inches(2.5), y + Inches(0.12), Inches(6.8), Inches(0.4))
        p2.text_frame.paragraphs[0].text = task
        p2.text_frame.paragraphs[0].font.size = Pt(12)
        p2.text_frame.paragraphs[0].font.color.rgb = TEXT
        y += Inches(0.68)
    ask = s12.shapes.add_textbox(Inches(0.55), Inches(4.85), Inches(8.9), Inches(0.55))
    ask.text_frame.paragraphs[0].text = "Discussion: prioritize vendor order now vs. wait for computational ranking of first subset?"
    ask.text_frame.paragraphs[0].font.size = Pt(13)
    ask.text_frame.paragraphs[0].font.italic = True
    ask.text_frame.paragraphs[0].font.color.rgb = CYAN_DARK
    add_slide_transition(s12, "push")

    # --- Slide 13: Thank you ---
    s13 = prs.slides.add_slide(blank)
    set_bg(s13, NAVY)
    if BANNER.exists():
        s13.shapes.add_picture(str(BANNER), Inches(0), Inches(0), width=Inches(10), height=Inches(2.0))
    t13 = s13.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(8.8), Inches(0.8))
    t13.text_frame.paragraphs[0].text = "Thank You"
    t13.text_frame.paragraphs[0].font.size = Pt(40)
    t13.text_frame.paragraphs[0].font.bold = True
    t13.text_frame.paragraphs[0].font.color.rgb = WHITE
    q = s13.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(8.8), Inches(1.5))
    qtf = q.text_frame
    for i, line in enumerate(
        [
            "Questions & feedback welcome",
            "Repository: github.com/jackyissocute/LiSPER",
            "LiSPER — Can engineered peptide ensembles make lithium recovery more selective?",
        ]
    ):
        p = qtf.paragraphs[0] if i == 0 else qtf.add_paragraph()
        p.text = line
        p.font.size = Pt(14 if i else 18)
        p.font.color.rgb = CYAN if i == 0 else SLATE
        p.font.bold = i == 0
    add_slide_transition(s13, "fade")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
